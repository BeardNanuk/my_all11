import numpy as np 

### plot t_x signals using numpy arrays

### created on Sun Mar 31 22:28:36 UTC 2019 
### created by Jiaze He  


import numpy as np
import matplotlib.pyplot as plt 



def _cscale(v, clip=100):
    """ Return limits for colormap.
    """
    perc = clip / 100.
    return -perc * abs(v).max(), perc * abs(v).max()
def _get_time(stream):
    """ Get fixed time vector for stream object.
    """
    dt = stream[0].stats.delta
    nt = len(stream[0].data)
    return np.arange(0, nt*dt, dt)

def _get_time2(dt,nt):
    """ Get fixed time vector for stream object.
    """
    #dt = stream[0].stats.delta
    #nt = len(stream[0].data)
    return np.arange(0, nt*dt, dt)

def get_regular_ticks(v, interval):
    """ Returns regular tick intervals.
    """
    from scipy.interpolate import interp1d

    f = interp1d(v, range(len(v)))
    begin = int(v[0] / interval) * interval
    end = v[-1]
    tick_labels = np.arange(begin, end, interval)
    ticks = f(tick_labels)

    return ticks, tick_labels

def getcoords(dir,num=0):
    from seisflows.tools.tools import exists
    from seisflows.plugins.solver_io.fortran_binary import _read
    for names in getnames(dir,num):
        if exists(names):
            x, z = _read(names[0]), _read(names[1])
            return x, z
    else:
        print('Coordinate files not found.')
        sys.exit()

def getargs():
## used in zsavesismo.py [+]  
    import argparse,os,sys
    parser = argparse.ArgumentParser()

    parser.add_argument('--workdir', nargs='?',
        default=os.getcwd())

    parser.add_argument('--parameter_file', nargs='?',
        default='parameters.py')

    parser.add_argument('--path_file', nargs='?',
        default='paths.py')

    return parser.parse_args()




def meshplot(x, y, z): 
    """ Plots values on 2D unstructured mesh
    """
    r = (max(x) - min(x))/(max(y) - min(y))
    rx = r/np.sqrt(1 + r**2)
    ry = 1/np.sqrt(1 + r**2)

    f = plt.figure(figsize=(10*rx, 10*ry))
    p = plt.tricontourf(x, y, z, 125, extend="both")
    plt.axis('image')
    return f, p

def getnames(dir,num=0):

    from os.path import abspath, join
    return [
        (abspath(join(dir, subdir, 'proc' + str(num).zfill(6) + '_x.bin')),
         abspath(join(dir, subdir, 'proc' + str(num).zfill(6) + '_z.bin'))) 
         for subdir in [
         '.',
         '../model', 
         '../output/model_init',
         '../../output/model_init',
         '../../../output/model_init',
         '../../../../output/model_init',
         ]]



def _get_offsets(stream):
    """ Return offsets.
    """
    nr = len(stream)
    offsets = np.zeros(nr)
    scalco = stream[0].stats.su.trace_header.scalar_to_be_applied_to_all_coordinates

    # set scale to km
    if scalco == 0:
        scalco = 1e-3 # assume coords are in m
    else:
        scalco = 1.0e-3 / scalco

    for i, tr in enumerate(stream):
        offsets[i] = (tr.stats.su.trace_header.group_coordinate_x -
                      tr.stats.su.trace_header.source_coordinate_x) * scalco
    return offsets



def textplot(struct_variable,filename):

    attributes = [a for a in dir(struct_variable)
                  if not (a.startswith('__') and a.endswith('__'))
                  and not (a == 'kernel')]

    fig = plt.figure(figsize=(10, 7))
    #fig.suptitle('bold figure suptitle', fontsize=14, fontweight='bold')

    ax = fig.add_subplot(111)
    fig.subplots_adjust(top=0.85)
    ax.set_title(str('parameters'))

    ax.set_xlabel('xlabel')
    ax.set_ylabel('ylabel')

    spc = 0.26

    for i in range(len(attributes)):
        offset = spc * (i)
        str_temp = attributes[i] + ' = ' + str((getattr(struct_variable, attributes[i])))
        #str_temp = attributes[i] + ' = ' + str("%.5f" % (getattr(struct_variable, attributes[i])))
       ###ax.text(0.3, 9.7-offset,textwrap.shorten(str_temp, 110), fontsize=11)
        ax.text(0.3, 9.7-offset, str_temp, fontsize=9)

    ax.axis([0, 10, 0, 10])
    #plt.axis('off')

    plt.savefig(filename,format='png', dpi=500)
    #plt.show()
    #plt.close()



def add_slide_ze(img_path,filename_pptx,left_start=None,top_start=None,width=None,height=None):

    from pptx import Presentation   
    from pptx.util import Inches 

    if left_start is None:
       left_start = 0

    if top_start is None:
       top_start = 0.9

    if width is None:
       width = 9

    if height is None:
       height = 5.7

    prs = Presentation(filename_pptx)

    blank_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = Inches(left_start)
    top = Inches(top_start)
    width = Inches(width)
    height = Inches(height)


    pic = slide.shapes.add_picture(img_path, left, top,width=width,height=height)

    # left = Inches(5)
    # height = Inches(5.5)
    # pic = slide.shapes.add_picture(img_path, left, top, height=height)

    prs.save(filename_pptx)






def plot_section_steps(data=np.random.rand(11, 11), clip=100, title='', dt = 5e-8, \
                       y_star = 0, y_end = 11, x_interval=1.0, y_interval=100.0, \
                       fraction_val=0.06, pad_val=0.04,cmap='seismic',\
                       fsize = 6,scale_factor = 1.5,fig_filename='',flag_savepng=0):

    nr = len(data[0,:])
    nt = len(data[1,:])

    d_aspect = (nr) / float(y_end-y_star)
    y_range=np.arange(y_star, y_end)
    
    fig, ax = plt.subplots(figsize=(fsize, scale_factor*fsize))

    im = ax.imshow(data[y_star:y_end,:], aspect=scale_factor*d_aspect, clim=_cscale(data, clip=clip))
    
    im.set_cmap(cmap)
    pad=pad_val
    fraction=fraction_val
    fig.colorbar(im,fraction=fraction, pad=pad, extend='both')
    ax.set_title(title)
    ax.set_xlabel('Receiver number')
    ax.set_ylabel(r'Time Steps')
    
    yticks, ytick_labels = get_regular_ticks(y_range, y_interval)
    ax.set_yticks(yticks)
    ax.set_yticklabels(ytick_labels)
    if flag_savepng is 1:
       plt.savefig(fig_filename,format='png', bbox_inches='tight') 
    #plt.close()
    fig.show()

def plot_one_signal_comp(data_syn=np.random.rand(11, 11),data_obs=np.random.rand(11, 11),\
                     dt_syn=0,dt_obs=0,source_num=0,receiver_num=0,iter_num=0,\
                    t_star_showsyn=0, t_end_showsyn=0,t_star_showobs=0, t_end_showobs=0,\
                    flag_range=0,flag_same_range=0,fig_filename='',flag_savepng=0):
    
    
    nt_syn=len(data_syn[:,0])
    nt_obs=len(data_obs[:,0])
    t_total_syn=np.arange(dt_syn,dt_syn*(nt_syn+1),dt_syn)
    t_total_obs=np.arange(dt_obs,dt_obs*(nt_obs+1),dt_obs)
  
    if flag_range is 0: 
        # show the total time step range
        t_end_showsyn=len(data_syn[:,0])
        t_end_showobs=len(data_obs[:,0])
    elif flag_same_range is 1: 
        # if show the same time range
        # for well-aligned data only
        t_star_showobs=int(t_total_syn[t_star_showsyn]/dt_obs)
        t_end_showobs=int(t_total_syn[t_end_showsyn]/dt_obs)
    
    fig, (ax2,ax3,ax5) = plt.subplots(nrows=3)
    ax2.plot(t_total_syn[t_star_showsyn:t_end_showsyn],data_syn[t_star_showsyn:t_end_showsyn,receiver_num],'g-')
    ax2.set_title('synthetic -source '+str(source_num)+' trace ' + str(receiver_num))
    ax2.set_xlabel(r'time ($\mu s$)')

    ax3.plot(t_total_obs[t_star_showobs:t_end_showobs],data_obs[t_star_showobs:t_end_showobs,receiver_num],'m-')
    ax3.set_title('observed -source '+str(source_num)+' trace ' + str(receiver_num))
    ax3.set_xlabel(r'time ($\mu s$)')

    ax5.plot(t_total_syn[t_star_showsyn:t_end_showsyn],\
             data_syn[t_star_showsyn:t_end_showsyn,receiver_num]/max(data_syn[t_star_showsyn:t_end_showsyn,receiver_num]),\
             'g-',label='syn')
    ax5.plot(t_total_obs[t_star_showobs:t_end_showobs],\
             data_obs[t_star_showobs:t_end_showobs,receiver_num]/max(data_obs[t_star_showobs:t_end_showobs,receiver_num]),\
             'm-',label='obs')
    ax5.set_title('together - iter num' + str(iter_num) + ' , source '+str(source_num)+' trace ' + str(receiver_num))
    ax5.set_xlabel(r'time ($\mu s$)')

    plt.tight_layout(rect=[0, 0, 1.5, 1.3])
    if flag_savepng is 1:
        plt.savefig(fig_filename,format='png', bbox_inches='tight')
    #plt.close(fig)


def plot_one_signal_comp2(data_syn=np.random.rand(11, 11),data_obs=np.random.rand(11, 11),\
                     data_adj=np.random.rand(11, 11),data_syn_init=np.random.rand(11, 11),\
                     dt=0,source_num=0,receiver_num=0,iter_num=0,\
                    t_star_showsyn=0, t_end_showsyn=0,t_star_showobs=0, t_end_showobs=0,\
                    t_star_showadj=0, t_end_showadj=0,\
                    flag_range=0,flag_same_range=0,fig_filename='',flag_savepng=0):
    
    nt=len(data_syn[:,0])
    t_total=np.arange(dt,dt*(nt+1),dt)

    t_star_showsyn_init=t_star_showsyn 
    t_end_showsyn_init=t_end_showsyn 

    if flag_range is 0:  
        # show the total time step range
        t_end_showsyn=len(data_syn[:,0])
        t_end_showobs=len(data_obs[:,0])
    elif flag_same_range is 1:  
        # if show the same time range
        # for well-aligned data only
        t_star_showobs=int(t_total[t_star_showsyn]/dt)
        t_end_showobs=int(t_total[t_end_showsyn]/dt)
    
    from seisflows.plugins import misfit
    ###### temp 
    from numpy.linalg import norm
    stf_load = np.loadtxt("obf/input/stf_ricker_dt40ns_dlst200_f300000.txt")
    stf_load_norm = - stf_load[:,1]/norm(stf_load[:,1],axis=0) 
    print('shape of stf_load:', stf_load.shape)
    print('shape of stf_load_norm:', stf_load_norm.shape)
    print('shape of data_obs_one_trace', data_obs[t_star_showobs:t_end_showobs,receiver_num].shape)

    fig, (ax1, ax2,ax3,ax5) = plt.subplots(nrows=4)
    ax1.plot(t_total[t_star_showsyn_init:t_end_showsyn_init],data_syn_init[t_star_showsyn:t_end_showsyn,receiver_num],'k-')
    ax1.plot(t_total[t_star_showobs:t_end_showobs],data_obs[t_star_showobs:t_end_showobs,receiver_num],'m-')
    ax1.set_title('original synthetic and observed - source '+str(source_num)+' trace ' + str(receiver_num))
    ax1.set_xlabel(r'time ($\mu s$)')

    ax2.plot(t_total[t_star_showsyn:t_end_showsyn],data_syn[t_star_showsyn:t_end_showsyn,receiver_num],'b-')
    ax2.plot(t_total[t_star_showsyn:t_end_showsyn],data_obs[t_star_showsyn:t_end_showsyn,receiver_num],'m-')
    ax2.set_title('current synthetic - source '+str(source_num)+' trace ' + str(receiver_num))
    ax2.set_xlabel(r'time ($\mu s$)')
#
    ax3.plot(t_total[t_star_showobs:t_end_showobs],data_obs[t_star_showobs:t_end_showobs,receiver_num],'m-')
    ax3.plot(t_total[t_star_showobs:t_end_showobs],stf_load_norm,'y-')
    ax3.set_title('observed - source '+str(source_num)+' trace ' + str(receiver_num))
    ax3.set_xlabel(r'time ($\mu s$)')

    #tt_cr = misfit.Traveltime(data_syn[:,receiver_num],data_obs[:,receiver_num],nt,dt)
    #print('tt_cr:',str(tt_cr))
    ax5.plot(t_total[t_star_showsyn:t_end_showsyn],data_adj[t_star_showadj:t_end_showadj,receiver_num],'g-')
    ax5.set_title('adjoint - iter num ' + str(iter_num) + ' source '+str(source_num)+' trace ' + str(receiver_num))
    ax5.set_xlabel(r'time ($\mu s$)')

    plt.tight_layout(rect=[0, 0, 2.8, 1.9])
    if flag_savepng is 1:
        plt.savefig(fig_filename,format='png', bbox_inches='tight')
    #plt.close(fig)





def plot_seismo_su(clip=100, title='', \
                       y_star = 0, y_end = 11, x_interval=1.0, y_interval=100.0, \
                       fig_filename='',flag_savepng=0,idir='',t_total=0,NREC=0 ):

    import prusct.io.read_data as rd
    from my_sfs.ze_plot_section import plot_section_steps

    data2dnp=rd.readebsu(idir=idir,t_total = t_total,NREC=NREC)
    plot_section_steps(data=data2dnp, clip=clip, y_star = y_star, y_end=y_end, title=title, y_interval=y_interval, fig_filename=fig_filename,flag_savepng=flag_savepng)



def seismofold_src(iter_num=0, srcnum=0, seismofold='', flag_data_or_scratch=0):
    import os
    srcnum_format= '%06d' % (int(srcnum))
    iter_num_format= 'traces_%04d' % (int(iter_num))

    if seismofold=='syn':
        sufile='Up_file_single.su'
    elif seismofold=='obs': 
        sufile='Up_file_single.su'
    elif seismofold=='adj': 
        sufile='Up_file_single.su.adj'
    
    #print('sourcenum:', srcnum)
    #print('srcnum_format:', srcnum_format)
    
    if flag_data_or_scratch==2:
         ipath=os.getcwd()+'/scratch/solver/' + srcnum_format + '/traces/' + seismofold + '/' + sufile
         syn_folder= '_scratch'
    elif flag_data_or_scratch==1:
         ipath=os.getcwd()+'/DATA/' + seismofold + '/' +  srcnum_format + '/' + sufile
         syn_folder= '_init'
    elif flag_data_or_scratch==3:
         ipath=os.getcwd()+'/output/' + iter_num_format + '/' + seismofold + '/' + srcnum_format +'/'+ sufile
         syn_folder= '_cont'

    return (sufile,srcnum_format,iter_num_format,syn_folder,ipath)




